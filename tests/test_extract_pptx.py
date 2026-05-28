from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SCRIPT = ROOT / "scripts" / "extract_pptx.py"


def write_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Unicode Test"
    slide.placeholders[1].text = "• bullet item\n中文内容"
    prs.save(path)


class ExtractPptxTests(unittest.TestCase):
    def test_default_stdout_is_summary_and_print_json_is_opt_in(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            pptx = tmp_path / "deck.pptx"
            default_out = tmp_path / "default"
            json_out = tmp_path / "json"
            write_pptx(pptx)

            default_result = subprocess.run(
                [sys.executable, str(SCRIPT), "--input", str(pptx), "--out-dir", str(default_out)],
                cwd=ROOT,
                capture_output=True,
                text=True,
                encoding="utf-8",
                check=True,
            )

            self.assertIn("Intermediate JSON written to:", default_result.stdout)
            self.assertNotIn('"slides"', default_result.stdout)

            json_result = subprocess.run(
                [
                    sys.executable,
                    str(SCRIPT),
                    "--input",
                    str(pptx),
                    "--out-dir",
                    str(json_out),
                    "--print-json",
                ],
                cwd=ROOT,
                capture_output=True,
                text=True,
                encoding="utf-8",
                check=True,
            )
            data = json.loads(json_result.stdout)
            self.assertEqual(data["slides"][0]["title"], "Unicode Test")
            self.assertTrue(any(ord(ch) > 127 for ch in data["slides"][0]["text"]))


if __name__ == "__main__":
    unittest.main()
