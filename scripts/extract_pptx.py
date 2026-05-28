#!/usr/bin/env python3
"""Extract a `.pptx` into the shared intermediate JSON plus image assets.

Usage:
    python extract_pptx.py --input <file.pptx> --out-dir <work-dir> [--assets-subdir assets] [--quiet] [--print-json]

Behavior:
    - One output item per slide
    - Text is collected from text frames, grouped shapes, tables, and chart labels
    - Title prefers slide.shapes.title, then falls back to the first short line
    - Notes come from the speaker notes area when available
    - Images are written to the assets directory and get a rough position label
    - The intermediate JSON is written to <out-dir>/intermediate.json
    - A compact extraction summary is printed to stdout by default
    - Full JSON is printed only when --print-json is passed

Exit codes:
    0  success
    1  argument error
    2  missing python-pptx
    3  file read or parse failure
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

from _schema import ImageRef, Intermediate, Slide, Source  # noqa: E402

# 1 inch = 914400 EMU
EMU_PER_INCH = 914400


def configure_utf8_stdio() -> None:
    """Avoid Windows GBK console crashes when output contains Unicode text."""
    for stream in (sys.stdout, sys.stderr):
        reconfigure = getattr(stream, "reconfigure", None)
        if reconfigure is None:
            continue
        try:
            reconfigure(encoding="utf-8")
        except Exception:
            pass


def iter_shapes(shapes):
    """Yield shapes recursively so grouped content is not silently skipped."""
    for shape in shapes:
        yield shape
        try:
            nested = shape.shapes
        except Exception:
            nested = None
        if nested is not None:
            yield from iter_shapes(nested)


def clean_text(text: str) -> str:
    """Normalize noisy text-frame output while preserving line breaks."""
    lines = [line.strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def table_text(shape) -> str:
    """Extract visible cell text from a PowerPoint table shape."""
    try:
        if not shape.has_table:
            return ""
    except Exception:
        return ""

    rows: list[str] = []
    try:
        for row in shape.table.rows:
            cells = [clean_text(cell.text) for cell in row.cells]
            row_text = " | ".join(cell for cell in cells if cell)
            if row_text:
                rows.append(row_text)
    except Exception:
        return ""
    return "\n".join(rows).strip()


def chart_text(shape) -> str:
    """Extract chart title, category labels, and series names when available."""
    try:
        if not shape.has_chart:
            return ""
        chart = shape.chart
    except Exception:
        return ""

    parts: list[str] = []

    try:
        if chart.has_title and chart.chart_title.text_frame:
            title = clean_text(chart.chart_title.text_frame.text)
            if title:
                parts.append(f"Chart title: {title}")
    except Exception:
        pass

    try:
        for plot in chart.plots:
            try:
                categories = [str(c) for c in plot.categories if str(c)]
                if categories:
                    parts.append("Chart categories: " + ", ".join(categories[:20]))
            except Exception:
                pass

            try:
                series_names = []
                for series in plot.series:
                    name = str(getattr(series, "name", "") or "").strip()
                    if name:
                        series_names.append(name)
                if series_names:
                    parts.append("Chart series: " + ", ".join(series_names[:20]))
            except Exception:
                pass
    except Exception:
        pass

    return "\n".join(parts).strip()


def extract(input_path: Path, out_dir: Path, assets_subdir: str = "") -> Intermediate:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Emu  # noqa: F401
    except ImportError:
        print(
            "Missing dependency: python-pptx. Install it with: pip install python-pptx",
            file=sys.stderr,
        )
        sys.exit(2)

    if not input_path.exists():
        print(f"Input file does not exist: {input_path}", file=sys.stderr)
        sys.exit(3)

    try:
        prs = Presentation(str(input_path))
    except Exception as e:
        print(f"Failed to open PPTX: {e}", file=sys.stderr)
        sys.exit(3)

    assets_dir = out_dir / assets_subdir if assets_subdir else out_dir
    assets_dir.mkdir(parents=True, exist_ok=True)

    # Metadata title
    meta_title = ""
    try:
        meta_title = prs.core_properties.title or ""
    except Exception:
        pass
    if not meta_title:
        meta_title = input_path.stem

    page_w_emu = prs.slide_width or 0
    page_h_emu = prs.slide_height or 0

    inter = Intermediate(
        source=Source(
            path=str(input_path.resolve()),
            format="pptx",
            title=meta_title,
            page_count=len(prs.slides),
        ),
        slides=[],
    )

    for slide_index, slide in enumerate(prs.slides, start=1):
        # ----- title -----
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text.strip()
        except Exception:
            title = ""

        # ----- text -----
        text_parts: list[str] = []
        seen_text: set[str] = set()
        for shape in iter_shapes(slide.shapes):
            try:
                if shape.has_text_frame:
                    t = clean_text(shape.text_frame.text)
                    if t and t != title and t not in seen_text:
                        text_parts.append(t)
                        seen_text.add(t)
            except Exception:
                pass

            for extra_text in (table_text(shape), chart_text(shape)):
                if extra_text and extra_text not in seen_text:
                    text_parts.append(extra_text)
                    seen_text.add(extra_text)
        text = "\n".join(text_parts).strip()

        # If there is no explicit title, try the first line of the text.
        if not title and text:
            first_line = text.split("\n", 1)[0].strip()
            if 0 < len(first_line) <= 60:
                title = first_line

        # ----- notes -----
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = (
                    slide.notes_slide.notes_text_frame.text.strip()
                    if slide.notes_slide.notes_text_frame
                    else ""
                )
        except Exception:
            notes = ""

        # ----- images -----
        images: list[ImageRef] = []
        img_idx = 0
        for shape in iter_shapes(slide.shapes):
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
            except Exception:
                continue

            try:
                image = shape.image
                blob = image.blob
                ext = image.ext or "png"
            except Exception:
                continue

            # Rough size estimate (EMU -> px at ~96 dpi)
            try:
                w_emu = shape.width or 0
                h_emu = shape.height or 0
                width_px = int(w_emu / EMU_PER_INCH * 96) if w_emu else 0
                height_px = int(h_emu / EMU_PER_INCH * 96) if h_emu else 0
            except Exception:
                width_px = height_px = 0

            # Skip tiny images, which are usually decorative.
            if width_px and height_px and (width_px < 100 or height_px < 100):
                continue

            img_idx += 1
            img_id = f"slide{slide_index}_img{img_idx}"
            filename = f"{img_id}.{ext}"
            (assets_dir / filename).write_bytes(blob)
            rel_path = (
                f"{assets_subdir}/{filename}" if assets_subdir else filename
            )

            # Rough position estimate
            position = "inset"
            try:
                if page_w_emu and page_h_emu and w_emu and h_emu:
                    ratio = (w_emu * h_emu) / (page_w_emu * page_h_emu)
                    if ratio > 0.7:
                        position = "full"
                    elif ratio > 0.3:
                        position = "center"
                    else:
                        position = "inset"
            except Exception:
                pass

            images.append(
                ImageRef(
                    id=img_id,
                    path=rel_path,
                    width=width_px,
                    height=height_px,
                    position=position,
                    context_text=text[:300],
                )
            )

        inter.slides.append(
            Slide(
                index=slide_index,
                title=title,
                text=text,
                notes=notes,
                images=images,
            )
        )

    return inter


def main() -> int:
    configure_utf8_stdio()
    ap = argparse.ArgumentParser(description="Extract a .pptx into the ppt2notes intermediate JSON")
    ap.add_argument("--input", required=True, help="Input PPTX path")
    ap.add_argument("--out-dir", required=True, help="Working directory")
    ap.add_argument("--assets-subdir", default="", help="Optional subdirectory for extracted images")
    ap.add_argument("--quiet", action="store_true", help="Suppress the default summary output")
    ap.add_argument(
        "--print-json",
        action="store_true",
        help="Print the full intermediate JSON to stdout. By default only a summary is printed.",
    )
    args = ap.parse_args()

    input_path = Path(args.input).expanduser().resolve()
    out_dir = Path(args.out_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    inter = extract(input_path, out_dir, args.assets_subdir)

    json_path = out_dir / "intermediate.json"
    json_path.write_text(inter.to_json(), encoding="utf-8")

    n_imgs = sum(len(s.images) for s in inter.slides)
    if args.print_json:
        print(inter.to_json())
        if not args.quiet:
            print(f"Intermediate JSON written to: {json_path}", file=sys.stderr)
            print(f"Slides: {len(inter.slides)}", file=sys.stderr)
            print(f"Extracted images: {n_imgs}", file=sys.stderr)
    elif not args.quiet:
        print(f"Intermediate JSON written to: {json_path}")
        print(f"Slides: {len(inter.slides)}")
        print(f"Extracted images: {n_imgs}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
