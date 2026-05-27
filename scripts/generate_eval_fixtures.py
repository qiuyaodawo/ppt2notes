#!/usr/bin/env python3
"""Generate synthetic fixtures for the ppt2notes eval cases."""
from __future__ import annotations

import argparse
import sys
from pathlib import Path


def require_deps():
    try:
        import fitz
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        print(
            "Missing dependency. Install fixture dependencies with: "
            "pip install pymupdf python-pptx",
            file=sys.stderr,
        )
        raise SystemExit(2) from exc
    return fitz, Presentation, Inches


def add_text_slide(prs, title: str, bullets: list[str], notes: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, item in enumerate(bullets):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        para.text = item
        para.level = 0
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def write_concept_heavy_zh(out_dir: Path, Presentation) -> None:
    prs = Presentation()
    slides = [
        ("机器学习导论", ["监督学习", "训练集与测试集", "泛化能力"], "强调泛化不是记住训练样本。"),
        ("线性模型", ["线性回归", "权重与偏置", "均方误差"], "补充最小二乘的直觉。"),
        ("分类问题", ["逻辑回归", "决策边界", "交叉熵损失"], "解释概率输出和阈值。"),
        ("优化方法", ["梯度下降", "学习率", "局部最小值"], "提醒学习率过大可能震荡。"),
        ("过拟合", ["训练误差下降", "验证误差上升", "正则化"], "联系模型复杂度。"),
        ("模型评估", ["准确率", "召回率", "F1 分数"], "说明类别不平衡时准确率会误导。"),
        ("课程总结", ["从数据到模型", "从训练到评估", "从误差到改进"], ""),
    ]
    for title, bullets, notes in slides:
        add_text_slide(prs, title, bullets, notes)
    prs.save(out_dir / "concept_heavy_zh.pptx")


def write_concept_followup_zh(out_dir: Path, Presentation) -> None:
    prs = Presentation()
    slides = [
        (
            "从线性模型到神经网络",
            ["上一讲回顾：线性模型使用权重和偏置", "本讲目标：引入非线性表示", "多层结构提升表达能力"],
            "承接上一讲的监督学习、损失函数和梯度下降。",
        ),
        (
            "隐藏层",
            ["隐藏层学习中间表示", "激活函数引入非线性", "层数影响表达能力"],
            "提醒学生不要把隐藏层理解成黑箱魔法。",
        ),
        (
            "训练神经网络",
            ["损失函数仍然衡量预测误差", "梯度下降仍然更新参数", "反向传播用于高效求梯度"],
            "强调这是上一讲优化思想的延伸。",
        ),
        (
            "与过拟合的关系",
            ["模型更复杂时更容易过拟合", "需要验证集", "可用正则化和早停"],
            "联系上一讲模型评估。",
        ),
        (
            "小结",
            ["线性模型是基础", "神经网络扩展了表示能力", "训练目标和评估逻辑保持一致"],
            "",
        ),
    ]
    for title, bullets, notes in slides:
        add_text_slide(prs, title, bullets, notes)
    prs.save(out_dir / "concept_followup_zh.pptx")


def render_png(fitz, path: Path, title: str, lines: list[str], width: int = 900, height: int = 520) -> None:
    doc = fitz.open()
    page = doc.new_page(width=width, height=height)
    page.draw_rect(fitz.Rect(30, 30, width - 30, height - 30), color=(0.1, 0.3, 0.6), width=3)
    page.insert_text((60, 90), title, fontsize=28, color=(0.05, 0.12, 0.2))
    y = 150
    for line in lines:
        page.insert_text((80, y), line, fontsize=21, color=(0.1, 0.1, 0.1))
        y += 44
    pix = page.get_pixmap(alpha=False)
    pix.save(str(path))
    doc.close()


def write_image_heavy(out_dir: Path, fitz, Presentation, Inches) -> None:
    diagram = out_dir / "_diagram.png"
    logo = out_dir / "_logo.png"
    render_png(
        fitz,
        diagram,
        "Three-stage pipeline",
        ["1. Collect data", "2. Train model", "3. Evaluate on validation set"],
    )
    render_png(fitz, logo, "COURSE", ["decorative badge"], width=240, height=180)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "模型训练流程"
    slide.shapes.add_picture(str(diagram), Inches(0.8), Inches(1.3), width=Inches(7.4))
    slide.shapes.add_picture(str(logo), Inches(8.0), Inches(0.2), width=Inches(1.5))

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "验证集的作用"
    slide2.placeholders[1].text = "验证集用于调参和早停，不应参与最终测试。"
    slide2.shapes.add_picture(str(logo), Inches(8.0), Inches(0.2), width=Inches(1.5))
    prs.save(out_dir / "image_heavy.pptx")

    diagram.unlink(missing_ok=True)
    logo.unlink(missing_ok=True)


def write_tiny_deck(out_dir: Path, Presentation) -> None:
    prs = Presentation()
    add_text_slide(prs, "递归", ["函数调用自身", "基例", "递归式"], "重点是基例必须收敛。")
    add_text_slide(prs, "例子", ["阶乘", "斐波那契", "树遍历"], "")
    add_text_slide(prs, "常见错误", ["缺少基例", "重复计算", "栈溢出"], "")
    prs.save(out_dir / "tiny_deck.pptx")


def write_math_heavy_pdf(out_dir: Path, fitz) -> None:
    formula = out_dir / "_formula.png"
    render_png(
        fitz,
        formula,
        "Gradient formula",
        ["dL/dW = dL/dy * dy/dz * dz/dW", "W <- W - eta * dL/dW"],
        width=760,
        height=320,
    )

    doc = fitz.open()
    page = doc.new_page(width=842, height=595)
    page.insert_text((72, 72), "Backpropagation and Optimization", fontsize=22)
    page.insert_text(
        (72, 120),
        "A neural network minimizes a loss function by repeatedly applying the chain rule.",
        fontsize=14,
    )
    page.insert_image(fitz.Rect(90, 180, 750, 460), filename=str(formula))

    page2 = doc.new_page(width=842, height=595)
    page2.insert_text((72, 72), "Regularization", fontsize=22)
    page2.insert_text(
        (72, 120),
        "L2 regularization adds lambda times the squared weight norm to the objective.",
        fontsize=14,
    )
    page2.insert_text((72, 170), "J(W) = L(W) + lambda * ||W||^2", fontsize=16)
    doc.save(out_dir / "math_heavy_en.pdf")
    doc.close()
    formula.unlink(missing_ok=True)


def write_scanned_pdf(out_dir: Path, fitz) -> None:
    png = out_dir / "_scanned_page.png"
    render_png(
        fitz,
        png,
        "Scanned lecture page",
        ["This text is rendered into an image.", "The output PDF has no text layer."],
        width=842,
        height=595,
    )

    doc = fitz.open()
    for _ in range(5):
        page = doc.new_page(width=842, height=595)
        page.insert_image(fitz.Rect(0, 0, 842, 595), filename=str(png))
    doc.save(out_dir / "scanned.pdf")
    doc.close()
    png.unlink(missing_ok=True)


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate ppt2notes eval fixtures")
    parser.add_argument("--out-dir", default="evals/fixtures", help="Fixture output directory")
    args = parser.parse_args()

    fitz, Presentation, Inches = require_deps()
    out_dir = Path(args.out_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    write_concept_heavy_zh(out_dir, Presentation)
    write_concept_followup_zh(out_dir, Presentation)
    write_image_heavy(out_dir, fitz, Presentation, Inches)
    write_tiny_deck(out_dir, Presentation)
    write_math_heavy_pdf(out_dir, fitz)
    write_scanned_pdf(out_dir, fitz)

    print(f"Generated eval fixtures in: {out_dir}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
